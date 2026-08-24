"""PPT 生成服务 — python-pptx 精美幻灯片 + 可选 TTS 配音"""
import os, uuid, tempfile
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import httpx

EXPORT_DIR = os.path.join(os.path.dirname(__file__), "..", "uploads", "export")
os.makedirs(EXPORT_DIR, exist_ok=True)

# ══════════════════════════════════════════════
# 配色方案
# ══════════════════════════════════════════════
THEMES = {
    "dark": {
        "bg":        RGBColor(0x0F, 0x17, 0x2A),
        "bg_light":  RGBColor(0x1E, 0x29, 0x3B),
        "text":      RGBColor(0xE2, 0xE8, 0xF0),
        "accent":    RGBColor(0x63, 0x66, 0xF1),
        "accent2":   RGBColor(0xA5, 0xB4, 0xFC),
        "code_bg":   RGBColor(0x0B, 0x0F, 0x19),
        "code_text": RGBColor(0xC7, 0xD2, 0xFE),
    },
    "tech": {
        "bg":        RGBColor(0x0B, 0x0F, 0x19),
        "bg_light":  RGBColor(0x16, 0x1B, 0x22),
        "text":      RGBColor(0xC9, 0xD1, 0xD9),
        "accent":    RGBColor(0x58, 0xA6, 0xFF),
        "accent2":   RGBColor(0x79, 0xC0, 0xFF),
        "code_bg":   RGBColor(0x01, 0x04, 0x09),
        "code_text": RGBColor(0xE6, 0xED, 0xF3),
    },
    "warm": {
        "bg":        RGBColor(0xFA, 0xFA, 0xFA),
        "bg_light":  RGBColor(0xF0, 0xF0, 0xF0),
        "text":      RGBColor(0x1F, 0x29, 0x37),
        "accent":    RGBColor(0xF5, 0x9E, 0x0B),
        "accent2":   RGBColor(0xDB, 0x27, 0x00),
        "code_bg":   RGBColor(0x1E, 0x1E, 0x1E),
        "code_text": RGBColor(0xF8, 0xF8, 0xF2),
    },
}

_W = Inches(13.333)  # 16:9
_H = Inches(7.5)
_FONT_CN = "Microsoft YaHei"
_FONT_EN = "Consolas"


def _add_bg(slide, color):
    """设置纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, color):
    """添加矩形色块"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # 无边框
    return shape


def _text_box(slide, left, top, width, height, text, font_size=18, color=None, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or RGBColor(0xE2, 0xE8, 0xF0)
    p.font.bold = bold
    p.font.name = font_name or _FONT_CN
    p.alignment = alignment
    return txBox


def _slide_title(slide, text, theme, subtitle=None):
    """标题页"""
    t = THEMES[theme]
    _add_bg(slide, t["bg"])
    # 装饰色块
    _add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), t["accent"])
    # 标题
    _text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
              text, font_size=40, color=t["text"], bold=True, alignment=PP_ALIGN.LEFT)
    # 副标题
    if subtitle:
        _text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
                  subtitle, font_size=18, color=t["accent2"], alignment=PP_ALIGN.LEFT)
    # 页脚
    _text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
              "Ray的垃圾站 · AI 生成", font_size=10, color=t["bg_light"], alignment=PP_ALIGN.LEFT)


def _slide_section(slide, text, theme):
    """分隔页"""
    t = THEMES[theme]
    _add_bg(slide, t["accent"])
    _text_box(slide, Inches(2), Inches(2.8), Inches(9), Inches(2),
              text, font_size=44, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)


async def _download_image(url: str) -> bytes | None:
    """下载远程图片"""
    try:
        async with httpx.AsyncClient(timeout=15.0, follow_redirects=True) as c:
            r = await c.get(url, headers={"User-Agent": "Mozilla/5.0"})
            if r.status_code == 200 and len(r.content) > 500:
                return r.content
    except Exception:
        pass
    return None


def _slide_content(slide, title, theme, bullets, code_block=None, image_bytes=None):
    """正文页：标题 + 要点列表 + 可选图片/代码块"""
    t = THEMES[theme]
    _add_bg(slide, t["bg"])
    _add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.04), t["accent"])
    _text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
              title, font_size=28, color=t["text"], bold=True)
    _add_shape(slide, Inches(1), Inches(1.3), Inches(6), Inches(0.02), t["accent"])

    has_image = image_bytes and len(image_bytes) > 500

    # 要点（左侧或全宽）
    text_width = Inches(7) if has_image else Inches(11.5)
    y = Inches(1.6)
    for b in bullets[:6]:
        _text_box(slide, Inches(1.2), y, text_width, Inches(0.5),
                  f"  {b}", font_size=16, color=t["text"])
        y += Inches(0.45)

    # 配图（右侧）
    if has_image:
        try:
            img_stream = BytesIO(image_bytes)
            slide.shapes.add_picture(img_stream, Inches(8.5), Inches(1.6),
                                     Inches(4.2), Inches(4.5))
        except Exception:
            pass

    # 代码块（图片下方或独立区域）
    if code_block:
        code_top = Inches(6.2) if has_image else Inches(1.6)
        _add_shape(slide, Inches(8.3), code_top, Inches(4.5), Inches(1.3), t["code_bg"])
        code_box = slide.shapes.add_textbox(Inches(8.5), code_top + Inches(0.05),
                                            Inches(4.2), Inches(1.2))
        tf = code_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code_block[:300]
        p.font.size = Pt(9)
        p.font.color.rgb = t["code_text"]
        p.font.name = _FONT_EN


def _slide_ending(slide, text, theme):
    """结尾页"""
    t = THEMES[theme]
    _add_bg(slide, t["bg"])
    _add_shape(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06), t["accent"])
    _text_box(slide, Inches(2), Inches(2.2), Inches(9), Inches(1.2),
              text, font_size=36, color=t["text"], bold=True, alignment=PP_ALIGN.CENTER)
    _text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
              "Ray的垃圾站 · 感谢收看", font_size=14, color=t["accent2"], alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# TTS 配音（Edge TTS，可靠免费）
# ══════════════════════════════════════════════

async def _edge_tts_audio(text: str) -> bytes:
    """Edge TTS 生成 MP3"""
    import edge_tts
    out = os.path.join(tempfile.gettempdir(), f"tts_{uuid.uuid4().hex[:8]}.mp3")
    try:
        voice = "zh-CN-XiaoxiaoNeural"
        await edge_tts.Communicate(text, voice).save(out)
        with open(out, "rb") as f:
            return f.read()
    finally:
        try:
            os.remove(out)
        except Exception:
            pass


# ══════════════════════════════════════════════
# 主入口
# ══════════════════════════════════════════════

async def generate_presentation(script: dict) -> dict:
    """根据脚本生成 PPT，返回 {url, filename, slide_count}"""
    title = script.get("title", "演示文稿")
    theme = script.get("theme", "dark")
    if theme not in THEMES:
        theme = "dark"
    slides_data = script.get("slides", script.get("segments", []))
    if not slides_data:
        return {"error": "脚本中没有幻灯片内容"}

    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H

    for i, s in enumerate(slides_data):
        stype = s.get("type", "content")
        stext = s.get("title", s.get("text", ""))
        bullets = s.get("bullets", [])
        subtitle = s.get("subtitle", "")
        code = s.get("code", "")
        notes = s.get("notes", stext)
        image_url = s.get("image_url", "")

        # 下载配图
        image_bytes = None
        if image_url:
            image_bytes = await _download_image(image_url)

        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)

        if stype == "title" or i == 0:
            _slide_title(slide, stext, theme, subtitle)
        elif stype == "section":
            _slide_section(slide, stext, theme)
        elif stype == "ending":
            _slide_ending(slide, stext, theme)
        else:
            if not bullets and stext:
                bullets = [stext] if len(stext) > 20 else []
            _slide_content(slide, stext, theme, bullets, code, image_bytes)

        # 添加演讲者备注
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

    filename = f"slides_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(EXPORT_DIR, filename)
    prs.save(filepath)

    url = f"/uploads/export/{filename}"
    return {
        "status": "success",
        "message": f"PPT已生成：{filename}，共{len(slides_data)}张幻灯片",
        "url": url,
        "filename": filename,
        "slide_count": len(slides_data),
    }
