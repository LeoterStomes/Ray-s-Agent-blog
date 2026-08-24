"""Agent tools definition + executor"""
import json, re, uuid
from datetime import datetime, timezone
import httpx
from sqlalchemy.orm import Session
from sqlalchemy import or_
from models import KnowledgeArticle, KnowledgeCategory

TOOLS_SCHEMA = [
    {"type":"function","function":{"name":"search_articles","description":"Semantic search the entire knowledge base (blog articles + imported reference docs like PDFs). Returns relevant chunks with similarity scores. Use this for any knowledge lookup - one search covers everything.","parameters":{"type":"object","properties":{"query":{"type":"string","description":"Search query in natural language"},"limit":{"type":"integer","description":"Max results, default 5","default":5}},"required":["query"]}}},
    {"type":"function","function":{"name":"get_article","description":"Get full article content by ID","parameters":{"type":"object","properties":{"article_id":{"type":"string","description":"Article ID (UUID)"}},"required":["article_id"]}}},
    {"type":"function","function":{"name":"get_categories","description":"List all blog categories with article counts","parameters":{"type":"object","properties":{}}}},
    {"type":"function","function":{"name":"recommend_articles","description":"Recommend related articles by category","parameters":{"type":"object","properties":{"article_id":{"type":"string","description":"Current article ID"},"limit":{"type":"integer","description":"Max results, default 3","default":3}},"required":["article_id"]}}},
    {"type":"function","function":{"name":"search_web","description":"Web search via DuckDuckGo→Bing→SearXNG multi-engine fallback. No JS required, returns text results.","parameters":{"type":"object","properties":{"query":{"type":"string","description":"Search keyword in natural language, e.g. 'github trending projects this week'"}},"required":["query"]}}},
    {"type":"function","function":{"name":"read_url","description":"Read any webpage URL and return text content","parameters":{"type":"object","properties":{"url":{"type":"string","description":"Full URL to read"}},"required":["url"]}}},
    {"type":"function","function":{"name":"get_recent_articles","description":"Get the most recent published articles","parameters":{"type":"object","properties":{"limit":{"type":"integer","description":"Max results, default 5","default":5}}}}},
    {"type":"function","function":{"name":"create_draft","description":"Create a draft blog article","parameters":{"type":"object","properties":{"title":{"type":"string","description":"Article title"},"content":{"type":"string","description":"Article content (HTML)"},"summary":{"type":"string","description":"Summary, optional"},"tags":{"type":"string","description":"Comma-separated tags"},"category_id":{"type":"integer","description":"Category ID, optional"}},"required":["title","content"]}}},
    {"type":"function","function":{"name":"make_mindmap","description":"Generate a mind map HTML file for embedding in blog articles. Returns a public URL and iframe code to insert into article content. Use this ONLY when writing/publishing articles, NOT for showing mind maps in chat.","parameters":{"type":"object","properties":{"markdown":{"type":"string","description":"Markdown outline for the mind map"}},"required":["markdown"]}}},
    {"type":"function","function":{"name":"read_document","description":"Read PDF/DOC/DOCX document content from URL","parameters":{"type":"object","properties":{"url":{"type":"string","description":"Document URL"}},"required":["url"]}}},
    {"type":"function","function":{"name":"summarize_url","description":"Summarize a webpage URL using AI. Returns a concise Chinese summary","parameters":{"type":"object","properties":{"url":{"type":"string","description":"Webpage URL to summarize"}},"required":["url"]}}},
    {"type":"function","function":{"name":"summarize_text","description":"Summarize text content using AI. Returns a concise Chinese summary","parameters":{"type":"object","properties":{"content":{"type":"string","description":"Text to summarize"},"max_length":{"type":"integer","description":"Max summary length, default 300","default":300}},"required":["content"]}}},
    {"type":"function","function":{"name":"export_file","description":"Export content as PDF/DOCX/TXT file. Returns download URL","parameters":{"type":"object","properties":{"title":{"type":"string","description":"Document title"},"content":{"type":"string","description":"Content to export (HTML or plain text)"},"format":{"type":"string","description":"File format: pdf / docx / txt","default":"pdf"}},"required":["title","content"]}}},
    {"type":"function","function":{"name":"generate_presentation","description":"Generate a professional PowerPoint presentation. Provide JSON with title, theme (dark/tech/warm), and slides. Each content slide can have title, bullets, image_url (for embedding pictures), code, and speaker notes. Use extract_images to get real project images from web pages.","parameters":{"type":"object","properties":{"script_json":{"type":"string","description":"JSON: {title, theme, slides: [{type, title, subtitle, bullets, image_url, code, notes}]}"}},"required":["script_json"]}}},
    {"type":"function","function":{"name":"extract_images","description":"Visit a webpage and extract all image URLs from it. Returns up to 10 image URLs with descriptions. Useful for getting real project logos, screenshots, and article images for PPTs.","parameters":{"type":"object","properties":{"url":{"type":"string","description":"Full webpage URL to extract images from"}},"required":["url"]}}},
    {"type":"function","function":{"name":"generate_weekly_video","description":"Generate a video from PPT slides + voiceover. First creates a PPT, then converts each slide to a video frame with TTS narration from speaker notes. Returns download URL for the MP4 video.","parameters":{"type":"object","properties":{"script_json":{"type":"string","description":"JSON: {title, theme, slides: [{type, title, subtitle, bullets, notes}]}. Same format as generate_presentation but notes field is used for TTS voiceover."}},"required":["script_json"]}}},
]

async def execute_tool(name: str, args: dict, db: Session) -> str:
    if name == "search_articles": return _search_articles(db, args.get("query",""), args.get("limit",5))
    elif name == "get_article": return _get_article(db, args.get("article_id",""))
    elif name == "get_categories": return _get_categories(db)
    elif name == "recommend_articles": return _recommend_articles(db, args.get("article_id",""), args.get("limit",3))
    elif name == "search_web": return await _search_web(args.get("query",""))
    elif name == "read_url": return await _read_url(args.get("url",""))
    elif name == "get_recent_articles": return _get_recent(db, args.get("limit",5))
    elif name == "create_draft": return _create_draft(db, args)
    elif name == "read_document": return await _read_document(args.get("url",""))
    elif name == "summarize_url": return await _summarize_url(args.get("url",""))
    elif name == "summarize_text": return await _summarize_text(args.get("content",""), args.get("max_length",300))
    elif name == "export_file": return _export_file(args)
    elif name == "make_mindmap": return _make_mindmap(args.get("markdown",""))
    elif name == "generate_presentation": return await _generate_ppt(args.get("script_json",""))
    elif name == "extract_images": return await _extract_images(args.get("url",""))
    elif name == "generate_weekly_video": return await _generate_video(args.get("script_json",""))
    return json.dumps({"error":f"Unknown tool: {name}"}, ensure_ascii=False)

# --- implementations ---

def _search_articles(db, query, limit=5):
    """全库语义搜索（博客文章 + 外部文档）+ LIKE 降级"""
    if not query: return json.dumps({"error":"Need keyword"}, ensure_ascii=False)
    limit = min(max(limit,1),10)

    # 尝试 RAG 语义搜索（全库：博客 + 外部文档）
    try:
        from services.rag_service import search_all as rag_search
        results = rag_search(query, limit)
        if results:
            articles = []
            docs = []
            for r in results:
                if "article_id" in r and r["article_id"]:
                    a = db.query(KnowledgeArticle).filter(KnowledgeArticle.id == r["article_id"]).first()
                    if a:
                        d = _a2d(a)
                        d["score"] = r.get("score", 0)
                        d["chunk_text"] = r.get("chunk_text", "")
                        d["source_type"] = "article"
                        articles.append(d)
                elif "filename" in r and r["filename"]:
                    docs.append({
                        "filename": r["filename"],
                        "chunk_text": r.get("chunk_text", ""),
                        "score": r.get("score", 0),
                        "source_type": "reference_doc",
                    })
            if articles or docs:
                return json.dumps({
                    "count": len(articles) + len(docs),
                    "articles": articles,
                    "reference_docs": docs,
                    "search_type": "semantic",
                }, ensure_ascii=False)
    except Exception:
        pass  # RAG 不可用时降级

    # LIKE 降级（仅博客文章）
    kw = f"%{query}%"
    arts = db.query(KnowledgeArticle).filter(KnowledgeArticle.status==1,
        or_(KnowledgeArticle.title.like(kw),KnowledgeArticle.summary.like(kw),
            KnowledgeArticle.tags.like(kw))
    ).order_by(KnowledgeArticle.published_at.desc()).limit(limit).all()
    return json.dumps({"count":len(arts),"articles":[_a2d(a) for a in arts],"reference_docs":[],"search_type":"keyword"}, ensure_ascii=False)

def _get_article(db, aid):
    if not aid: return json.dumps({"error":"Need ID"}, ensure_ascii=False)
    a = db.query(KnowledgeArticle).filter(KnowledgeArticle.id==aid).first()
    if not a: return json.dumps({"error":"Not found"}, ensure_ascii=False)
    d = _a2d(a); d["content"] = (a.content or "")[:8000]
    d["content_truncated"] = len(a.content or "") > 8000
    return json.dumps(d, ensure_ascii=False)

def _get_categories(db):
    cats = db.query(KnowledgeCategory).all()
    r = [{"id":c.id,"name":c.category_name,"description":c.description,
          "article_count":db.query(KnowledgeArticle).filter(
              KnowledgeArticle.category_id==c.id,KnowledgeArticle.status==1).count()} for c in cats]
    return json.dumps({"categories":r}, ensure_ascii=False)

def _recommend_articles(db, aid, limit=3):
    if not aid: return json.dumps({"error":"Need ID"}, ensure_ascii=False)
    a = db.query(KnowledgeArticle).filter(KnowledgeArticle.id==aid).first()
    if not a: return json.dumps({"error":"Not found"}, ensure_ascii=False)
    q = db.query(KnowledgeArticle).filter(KnowledgeArticle.status==1,KnowledgeArticle.id!=aid)
    if a.category_id: q = q.filter(KnowledgeArticle.category_id==a.category_id)
    arts = q.order_by(KnowledgeArticle.published_at.desc()).limit(limit).all()
    return json.dumps({"count":len(arts),"articles":[_a2d(x) for x in arts]}, ensure_ascii=False)

def _get_recent(db, limit=5):
    limit = min(max(limit,1),10)
    arts = db.query(KnowledgeArticle).filter(KnowledgeArticle.status==1
        ).order_by(KnowledgeArticle.published_at.desc()).limit(limit).all()
    return json.dumps({"count":len(arts),"articles":[_a2d(a) for a in arts]}, ensure_ascii=False)

def _create_draft(db, args):
    title = args.get("title","Untitled")
    content = args.get("content","")
    if not content: return json.dumps({"error":"Content required"}, ensure_ascii=False)
    try:
        from models import User
        author = db.query(User).filter(User.user_type==2).first()
        author_id = author.id if author else 1
    except: author_id = 1
    try:
        a = KnowledgeArticle(id=str(uuid.uuid4()),title=title,summary=args.get("summary","") or "",
            content=content,tags=args.get("tags","") or "",category_id=args.get("category_id") or None,
            author_id=author_id,status=0,published_at=datetime.now(timezone.utc).replace(tzinfo=None))
        db.add(a); db.commit()
        return json.dumps({"id":a.id,"title":title,"url":f"/blog/{a.id}","status":"draft"}, ensure_ascii=False)
    except Exception as e: return json.dumps({"error":str(e)}, ensure_ascii=False)

def _a2d(a):
    return {"id":a.id,"title":a.title,"summary":a.summary,
        "category":a.category.category_name if a.category else "Unknown",
        "tags":a.tags,"read_count":a.read_count or 0,
        "published_at":a.published_at.isoformat() if a.published_at else "","url":f"/blog/{a.id}"}

async def _extract_images(url: str) -> str:
    """提取网页中的图片 URL"""
    if not url:
        return json.dumps({"error": "Need URL"}, ensure_ascii=False)
    try:
        async with httpx.AsyncClient(timeout=15.0, follow_redirects=True) as c:
            r = await c.get(url, headers={"User-Agent": "Mozilla/5.0"})
            if r.status_code != 200:
                return json.dumps({"error": f"HTTP {r.status_code}"}, ensure_ascii=False)
            html = r.text

        images = []
        # 提取 og:image（文章主图，最优先）
        og_match = re.search(r'<meta[^>]+property="og:image"[^>]+content="([^"]+)"', html, re.I)
        if og_match:
            images.append({"url": og_match.group(1), "type": "og:image", "label": "文章主图"})

        # 提取所有 <img> 的 src
        img_matches = re.finditer(r'<img[^>]+src="([^"]+)"[^>]*(?:alt="([^"]*)")?', html, re.I)
        for m in img_matches:
            src = m.group(1)
            alt = m.group(2) or ""
            # 过滤小图标和占位符
            if any(s in src.lower() for s in ['avatar', 'icon', 'logo-small', '1x1', 'pixel', 'spacer', 'tracking']):
                continue
            # 补全相对路径
            if src.startswith("//"):
                src = "https:" + src
            elif src.startswith("/"):
                from urllib.parse import urljoin
                src = urljoin(url, src)
            if src.startswith("http") and len(src) < 2000:
                label = alt[:50] if alt else "配图"
                images.append({"url": src, "type": "img", "label": label})

        # 去重，最多 10 张
        seen = set()
        unique = []
        for img in images:
            if img["url"] not in seen:
                seen.add(img["url"])
                unique.append(img)
        unique = unique[:10]

        return json.dumps({"url": url, "image_count": len(unique), "images": unique}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


async def _read_url(url):
    if not url: return json.dumps({"error":"Need URL"}, ensure_ascii=False)
    try:
        async with httpx.AsyncClient(timeout=10.0, follow_redirects=True) as c:
            r = await c.get(url, headers={"User-Agent":"Mozilla/5.0"})
            if r.status_code==200:
                t = r.text
                t = re.sub(r'<script[^>]*>.*?</script>','',t,flags=re.DOTALL)
                t = re.sub(r'<style[^>]*>.*?</style>','',t,flags=re.DOTALL)
                t = re.sub(r'<[^>]+>',' ',t); t = re.sub(r'\s+',' ',t).strip()
                if len(t)>3000: t = t[:8000]+"...(truncated)"
                return json.dumps({"url":url,"content":t}, ensure_ascii=False)
    except Exception as e: return json.dumps({"error":str(e)}, ensure_ascii=False)
    return json.dumps({"error":"Failed"}, ensure_ascii=False)

async def _read_document(url):
    if not url: return json.dumps({"error":"Need URL"}, ensure_ascii=False)
    try:
        async with httpx.AsyncClient(timeout=20.0, follow_redirects=True) as c:
            r = await c.get(url, headers={"User-Agent":"Mozilla/5.0"})
            if r.status_code!=200: return json.dumps({"error":f"HTTP {r.status_code}"}, ensure_ascii=False)
            content, ct = r.content, r.headers.get("content-type","").lower()
        ul = url.lower(); text = ""
        if "pdf" in ct or ul.endswith(".pdf"):
            try:
                from PyPDF2 import PdfReader; import io
                pages = [p.extract_text() or "" for p in PdfReader(io.BytesIO(content)).pages[:10]]
                text = "\n".join(pages)
            except Exception as e: return json.dumps({"error":f"PDF: {e}"}, ensure_ascii=False)
        elif "docx" in ct or "word" in ct or ul.endswith(".docx"):
            try:
                from docx import Document; import io
                text = "\n".join(p.text for p in Document(io.BytesIO(content)).paragraphs if p.text.strip())
            except Exception as e: return json.dumps({"error":f"DOCX: {e}"}, ensure_ascii=False)
        elif "pptx" in ct or "presentation" in ct or ul.endswith(".pptx"):
            try:
                from pptx import Presentation; import io
                prs = Presentation(io.BytesIO(content))
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    parts = [f"--- Slide {i+1} ---"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            parts.append(shape.text_frame.text)
                    slides_text.append("\n".join(parts))
                text = "\n\n".join(slides_text[:20])  # 最多 20 页
            except Exception as e: return json.dumps({"error":f"PPTX: {e}"}, ensure_ascii=False)
        elif "text/plain" in ct: text = content.decode("utf-8",errors="ignore")
        else:
            t = content.decode("utf-8",errors="ignore")
            t = re.sub(r'<script[^>]*>.*?</script>','',t,flags=re.DOTALL)
            t = re.sub(r'<style[^>]*>.*?</style>','',t,flags=re.DOTALL)
            t = re.sub(r'<[^>]+>',' ',t); text = re.sub(r'\s+',' ',t).strip()
        if not text.strip(): return json.dumps({"error":"No text extracted"}, ensure_ascii=False)
        if len(text)>3000: text = text[:8000]+"...(truncated)"
        return json.dumps({"url":url,"type":ct or "unknown","content":text}, ensure_ascii=False)
    except Exception as e: return json.dumps({"error":str(e)}, ensure_ascii=False)

def _make_mindmap(markdown: str) -> str:
    """生成脑图 HTML 并复制到 /uploads/export/，返回公开 URL"""
    import shutil, uuid, os
    if not markdown:
        return json.dumps({"error": "请提供 Markdown 内容"}, ensure_ascii=False)
    try:
        # 调用 markmap MCP 生成 HTML
        from services.mcp_client import _servers
        markmap_server = _servers.get("markmap")
        if not markmap_server:
            # fallback: 手动调用 markmap-mcp-server
            from services.mcp_client import MCPServer
            markmap_server = MCPServer('markmap', 'C:/Users/Administrator/AppData/Roaming/npm/markmap-mcp-server.cmd', [])
            if not markmap_server.start():
                return json.dumps({"error": "markmap 服务未启动"}, ensure_ascii=False)

        result = markmap_server.call_tool('mcp_markmap__markdown_to_mindmap', {'markdown': markdown})
        d = json.loads(result)
        src = json.loads(d.get("result", "{}")).get("filePath", "")
        if not src or not os.path.exists(src):
            return json.dumps({"error": f"脑图生成失败，路径: {src}"}, ensure_ascii=False)

        # 复制到公开目录，并注入浅色主题样式
        export_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads", "export")
        os.makedirs(export_dir, exist_ok=True)
        filename = f"mindmap_{uuid.uuid4().hex[:8]}.html"
        dst = os.path.join(export_dir, filename)
        shutil.copy2(src, dst)
        # 注入浅色主题 CSS，覆盖 markmap 默认暗色背景
        with open(dst, "r", encoding="utf-8") as f:
            html = f.read()
        light_css = "<style>body{background:#fff!important;color:#1e293b!important}body .markmap-node text{fill:#334155!important}body .markmap-link{stroke:#94a3b8!important}svg{background:#fff!important}</style>"
        html = html.replace("</head>", light_css + "</head>", 1)
        with open(dst, "w", encoding="utf-8") as f:
            f.write(html)
        url = f"/uploads/export/{filename}"
        return json.dumps({"url": url, "filename": filename, "iframe": f'<iframe src="{url}" width="100%" height="500" style="border:none;border-radius:8px;margin:16px 0;background:#fff"></iframe>'}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


async def _generate_ppt(script_json: str) -> str:
    """生成 PPT 幻灯片"""
    if not script_json:
        return json.dumps({"error": "请提供幻灯片脚本 JSON"}, ensure_ascii=False)
    try:
        script = json.loads(script_json) if isinstance(script_json, str) else script_json
        from .slide_service import generate_presentation
        result = await generate_presentation(script)
        return json.dumps(result, ensure_ascii=False)
    except json.JSONDecodeError:
        return json.dumps({"error": "脚本 JSON 格式错误"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


async def _generate_video(script_json: str) -> str:
    """PPT → 视频（v2）"""
    if not script_json:
        return json.dumps({"error": "请提供脚本 JSON"}, ensure_ascii=False)
    try:
        script = json.loads(script_json) if isinstance(script_json, str) else script_json
        from .video_service import ppt_to_video
        result = await ppt_to_video(script)
        return json.dumps(result, ensure_ascii=False)
    except json.JSONDecodeError:
        return json.dumps({"error": "脚本 JSON 格式错误"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def _export_file(args: dict) -> str:
    """Export content as PDF/DOCX/TXT/HTML"""
    from .export_service import export_pdf, export_docx, export_txt, export_html
    title = args.get("title", "Untitled")
    content = args.get("content", "")
    fmt = args.get("format", "pdf").lower()
    try:
        if fmt == "docx": r = export_docx(title, content)
        elif fmt == "txt": r = export_txt(title, content)
        elif fmt == "html": r = export_html(title, content)
        else: r = export_pdf(title, content)
        return json.dumps({"title": title, "format": fmt, "url": r["url"], "filename": r["filename"]}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)

async def _summarize_url(url):
    if not url: return json.dumps({"error":"Need URL"}, ensure_ascii=False)
    from .summarize import summarize_url
    try:
        summary = await summarize_url(url)
        return json.dumps({"url":url,"summary":summary}, ensure_ascii=False)
    except Exception as e: return json.dumps({"error":str(e)}, ensure_ascii=False)

async def _summarize_text(content, max_length=300):
    if not content: return json.dumps({"error":"Need content"}, ensure_ascii=False)
    from .summarize import summarize_text
    try:
        summary = await summarize_text(content, max_length)
        return json.dumps({"summary":summary}, ensure_ascii=False)
    except Exception as e: return json.dumps({"error":str(e)}, ensure_ascii=False)

async def _search_web(query, engine="auto"):
    if not query: return json.dumps({"query":query,"results":"No query"}, ensure_ascii=False)
    import asyncio as _asyncio
    from config import BAIDU_API_KEY
    # 多引擎故障转移：百度AI搜索 → DuckDuckGo(lite) → Bing → SearXNG
    engines = []
    if BAIDU_API_KEY:
        engines.append(("baidu", "https://qianfan.baidubce.com/v2/ai_search/web_search", {
            "messages": [{"content": query, "role": "user"}],
            "search_recency_filter": "month",
            "resource_type_filter": [{"type": "web", "top_k": 8}],
        }))
    engines += [
        ("searxng", "https://search.sapti.me/search", {"q": query, "format": "json"}),
        ("duckduckgo", "https://lite.duckduckgo.com/lite/", {"q": query}),
    ]
    if engine != "auto":
        engines = [e for e in engines if e[0] == engine] or engines

    async def _try_one(eng_name, eng_url, eng_params):
        hdrs = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
                "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8", "Content-Type": "application/json"}
        async with httpx.AsyncClient(timeout=8.0, follow_redirects=True) as c:
            # 百度 API 用 POST + JSON
            if eng_name == "baidu":
                r = await c.post(eng_url, headers={**hdrs, "X-Appbuilder-Authorization": f"Bearer {BAIDU_API_KEY}"},
                                 json=eng_params)
            else:
                r = await c.get(eng_url, params=eng_params, headers=hdrs)
            if r.status_code != 200: return None
            # 百度返回 JSON（字段名: references）
            if eng_name == "baidu":
                try:
                    data = r.json()
                    refs = data.get("references", [])[:10]
                    lines = []
                    for item in refs:
                        title = item.get("title", "")
                        url = item.get("url", "")
                        snippet = item.get("content", "") or item.get("summary", "")
                        if title and url:
                            lines.append(f"- {title}\\n  {url}\\n  {snippet[:150]}")
                    if lines:
                        return json.dumps({"engine": "baidu", "query": query, "results": '\\n'.join(lines)[:2000]}, ensure_ascii=False)
                except Exception:
                    pass
                return None
            # SearXNG 返回 JSON
            if eng_name == "searxng":
                try:
                    data = r.json()
                    results = data.get("results", [])[:8]
                    lines = []
                    for item in results:
                        title = item.get("title", "")
                        url = item.get("url", "")
                        snippet = item.get("content", "") or item.get("snippet", "")
                        if title and url:
                            lines.append(f"- {title}\\n  {url}\\n  {snippet[:100]}")
                    if lines:
                        return json.dumps({"engine": eng_name, "query": query, "results": '\\n'.join(lines)[:1500]}, ensure_ascii=False)
                except Exception:
                    pass
                return None
            # DuckDuckGo / Bing：提取文本
            text = _extract_text(r.text)
            if text and len(text) > 50:
                return json.dumps({"engine": eng_name, "query": query, "results": text[:1500]}, ensure_ascii=False)
            return None

    for eng_name, eng_url, eng_params in engines:
        try:
            result = await _asyncio.wait_for(_try_one(eng_name, eng_url, eng_params), timeout=10.0)
            if result: return result
        except (_asyncio.TimeoutError, Exception):
            continue
    return json.dumps({"query": query, "results": "搜索超时或暂无结果，请稍后重试"}, ensure_ascii=False)


def _extract_text(html):
    """HTML → text with result extraction"""
    t = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
    t = re.sub(r'<style[^>]*>.*?</style>', '', t, flags=re.DOTALL)
    t = re.sub(r'<[^>]+>', ' ', t)
    t = re.sub(r'\s+', ' ', t).strip()
    results = []
    for m in re.finditer(r'<a[^>]*href="(https?://[^"]+)"[^>]*>([^<]{5,80})</a>', t):
        lt = re.sub(r'\s+', ' ', m.group(2)).strip()
        if lt and not any(s in lt.lower() for s in ['首页','下一页','登录','注册','广告','更多']):
            results.append(f"- {lt}\n  {m.group(1)}")
    return "\n".join(results[:8]) if results else t[:1500]
